"""
add_data.py  —  Universal document ingestion for the RAG vector database.

Supported file types:
  Documents  : .pdf  .docx  .pptx  .xlsx  .odt
  Text / Code: .txt  .md  .rst  .log  .py  .js  .ts  .java  .cpp  .c
               .cs  .go  .rb  .php  .swift  .kt  .sql  .sh  .bat  .ps1
               .r  .html  .htm  .xml  .yaml  .yml  .toml  .ini  .cfg
  Data       : .csv  .json  .jsonl

Usage:
  python add_data.py                          # GUI file picker (default)
  python add_data.py --folder C:/my_docs      # ingest all files from a folder
  python add_data.py file1.pdf file2.docx     # ingest specific files
  python add_data.py --help                   # show all options
"""

import os
import sys
import json
import logging
import argparse
import io

# Force UTF-8 output on Windows to avoid cp1252 encode errors
if sys.stdout.encoding and sys.stdout.encoding.lower() != 'utf-8':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')
from pathlib import Path

from dotenv import load_dotenv

load_dotenv()

# ── Logging setup ─────────────────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  %(levelname)s  %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger(__name__)


# ── Supported extensions ──────────────────────────────────────────────────────

TEXT_EXTENSIONS = {
    ".txt", ".md", ".rst", ".log", ".py", ".js", ".ts", ".java",
    ".cpp", ".c", ".cs", ".go", ".rb", ".php", ".swift", ".kt",
    ".sql", ".sh", ".bat", ".ps1", ".r", ".toml", ".ini", ".cfg",
    ".conf", ".env", ".yaml", ".yml"
}
HTML_EXTENSIONS  = {".html", ".htm", ".xml"}
DATA_EXTENSIONS  = {".csv", ".json", ".jsonl"}
PDF_EXTENSIONS   = {".pdf"}
WORD_EXTENSIONS  = {".docx"}
PPT_EXTENSIONS   = {".pptx"}
EXCEL_EXTENSIONS = {".xlsx", ".xls"}
ODT_EXTENSIONS   = {".odt"}

ALL_SUPPORTED = (
    TEXT_EXTENSIONS | HTML_EXTENSIONS | DATA_EXTENSIONS |
    PDF_EXTENSIONS | WORD_EXTENSIONS | PPT_EXTENSIONS |
    EXCEL_EXTENSIONS | ODT_EXTENSIONS
)


# ── Optional dependency loader ────────────────────────────────────────────────

def _try_import(package: str, install_name: str = ""):
    try:
        import importlib
        return importlib.import_module(package)
    except ImportError:
        pip_name = install_name or package
        log.warning(f"Package '{pip_name}' not installed. "
                    f"Run:  pip install {pip_name}")
        return None


# ── Text extractors  (one per file family) ────────────────────────────────────

def extract_plain_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def extract_html(path: Path) -> str:
    bs4 = _try_import("bs4", "beautifulsoup4")
    if not bs4:
        return extract_plain_text(path)       # fallback: raw HTML
    from bs4 import BeautifulSoup
    raw = path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(raw, "html.parser")
    return soup.get_text(separator="\n")


def extract_pdf(path: Path) -> str:
    pdfplumber = _try_import("pdfplumber")
    if not pdfplumber:
        return ""
    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                pages.append(t)
    return "\n".join(pages)


def extract_docx(path: Path) -> str:
    docx = _try_import("docx", "python-docx")
    if not docx:
        return ""
    import docx as _docx
    doc = _docx.Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_pptx(path: Path) -> str:
    pptx = _try_import("pptx", "python-pptx")
    if not pptx:
        return ""
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n".join(parts)


def extract_xlsx(path: Path) -> str:
    openpyxl = _try_import("openpyxl")
    if not openpyxl:
        return ""
    import openpyxl as xl
    wb = xl.load_workbook(str(path), read_only=True, data_only=True)
    rows = []
    for sheet in wb.worksheets:
        rows.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            line = "  ".join(str(c) for c in row if c is not None)
            if line.strip():
                rows.append(line)
    return "\n".join(rows)


def extract_csv(path: Path) -> str:
    import csv
    rows = []
    with path.open(encoding="utf-8", errors="ignore", newline="") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append("  ".join(row))
    return "\n".join(rows)


def extract_json(path: Path) -> str:
    try:
        data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        return extract_plain_text(path)


def extract_jsonl(path: Path) -> str:
    lines = []
    for raw in path.read_text(encoding="utf-8", errors="ignore").splitlines():
        raw = raw.strip()
        if not raw:
            continue
        try:
            lines.append(json.dumps(json.loads(raw), ensure_ascii=False))
        except json.JSONDecodeError:
            lines.append(raw)
    return "\n".join(lines)


def extract_odt(path: Path) -> str:
    odfpy = _try_import("odf.opendocument", "odfpy")
    if not odfpy:
        return ""
    from odf.opendocument import load as odf_load
    from odf import text as odf_text
    doc = odf_load(str(path))
    parts = []
    for el in doc.getElementsByType(odf_text.P):
        t = "".join(node.data for node in el.childNodes
                    if node.nodeType == node.TEXT_NODE)
        if t.strip():
            parts.append(t)
    return "\n".join(parts)


# ── Routing ───────────────────────────────────────────────────────────────────

def extract_text(path: Path) -> str:
    """Dispatch to the right extractor based on file extension."""
    ext = path.suffix.lower()

    if ext in TEXT_EXTENSIONS:   return extract_plain_text(path)
    if ext in HTML_EXTENSIONS:   return extract_html(path)
    if ext in PDF_EXTENSIONS:    return extract_pdf(path)
    if ext in WORD_EXTENSIONS:   return extract_docx(path)
    if ext in PPT_EXTENSIONS:    return extract_pptx(path)
    if ext in EXCEL_EXTENSIONS:  return extract_xlsx(path)
    if ext in ODT_EXTENSIONS:    return extract_odt(path)
    if ext == ".csv":            return extract_csv(path)
    if ext == ".json":           return extract_json(path)
    if ext == ".jsonl":          return extract_jsonl(path)

    # Unknown extension — try reading as UTF-8 text anyway
    log.warning(f"Unknown extension '{ext}' for {path.name}, trying as plain text.")
    try:
        return extract_plain_text(path)
    except Exception:
        return ""


# ── Chunking ──────────────────────────────────────────────────────────────────

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    chunks, start = [], 0
    text = text.strip()
    while start < len(text):
        chunks.append(text[start : start + chunk_size].strip())
        start += chunk_size - overlap
    return [c for c in chunks if c]


# ── Database ingestion ────────────────────────────────────────────────────────

def ingest_file(path: Path, cur, chunk_size: int, overlap: int) -> int:
    """Extract → chunk → embed → insert one file. Returns chunks inserted."""
    log.info(f"[FILE] {path.name}  [{path.suffix or 'no ext'}]")
    try:
        text = extract_text(path)
    except Exception as e:
        log.error(f"   Could not read {path.name}: {e}")
        return 0

    if not text.strip():
        log.warning(f"   No text extracted from {path.name} — skipping.")
        return 0

    chunks = chunk_text(text, chunk_size=chunk_size, overlap=overlap)
    log.info(f"   → {len(chunks)} chunk(s)")

    from app.embeddings import get_embedding

    inserted = 0
    for i, chunk in enumerate(chunks, 1):
        try:
            emb = get_embedding(chunk)
            emb_str = "[" + ",".join(str(v) for v in emb) + "]"
            cur.execute(
                "INSERT INTO documents (content, embedding) VALUES (%s, %s::vector)",
                (chunk, emb_str),
            )
            inserted += 1
        except Exception as e:
            log.error(f"   ✘ Chunk {i}/{len(chunks)} failed: {e}")

    log.info(f"   >> {inserted}/{len(chunks)} chunks ingested")
    return inserted


def run_ingestion(file_paths: list[Path], chunk_size: int, overlap: int) -> None:
    from app.db import get_connection, init_db

    init_db()
    conn = get_connection()
    cur  = conn.cursor()

    total_files = total_chunks = 0
    failed_files = []

    print()
    print("=" * 56)
    print(f"  Ingesting {len(file_paths)} file(s) ...")
    print("=" * 56)

    for fp in file_paths:
        if not fp.exists():
            log.error(f"File not found: {fp}")
            failed_files.append(str(fp))
            continue
        n = ingest_file(fp, cur, chunk_size, overlap)
        if n > 0:
            total_files  += 1
            total_chunks += n
        else:
            failed_files.append(fp.name)

    conn.commit()
    cur.close()
    conn.close()

    print()
    print("=" * 56)
    print("  DONE!")
    print(f"  Files processed : {total_files}")
    print(f"  Chunks inserted : {total_chunks}")
    if failed_files:
        print(f"  Failed / skipped: {', '.join(failed_files)}")
    print("=" * 56)
    print()


# ── GUI file picker ───────────────────────────────────────────────────────────

def open_file_picker() -> list[Path]:
    """Open a native OS file-picker; returns selected file paths."""
    try:
        import tkinter as tk
        from tkinter import filedialog
    except ImportError:
        log.error("tkinter is not available. Use  --folder  or pass file paths directly.")
        sys.exit(1)

    root = tk.Tk()
    root.withdraw()
    root.attributes("-topmost", True)

    ext_list = " ".join(f"*{e}" for e in sorted(ALL_SUPPORTED))
    filetypes = [
        ("All supported files", ext_list),
        ("PDF files",           "*.pdf"),
        ("Word documents",      "*.docx"),
        ("PowerPoint",          "*.pptx"),
        ("Excel / CSV",         "*.xlsx *.xls *.csv"),
        ("Text / Markdown",     "*.txt *.md *.rst *.log"),
        ("Code files",          "*.py *.js *.ts *.java *.cpp *.c *.cs *.go *.rb *.sql"),
        ("JSON / JSONL",        "*.json *.jsonl"),
        ("HTML / XML",          "*.html *.htm *.xml"),
        ("All files",           "*.*"),
    ]

    print("\n  Opening file picker — select one or more files ...\n")
    selected = filedialog.askopenfilenames(
        title="Select files to ingest into RAG database",
        filetypes=filetypes,
    )
    root.destroy()

    if not selected:
        print("  No files selected. Exiting.")
        sys.exit(0)

    return [Path(p) for p in selected]


# ── Folder scan ───────────────────────────────────────────────────────────────

def collect_from_folder(folder: str, recursive: bool = False) -> list[Path]:
    base = Path(folder)
    if not base.is_dir():
        log.error(f"Folder not found: {base}")
        sys.exit(1)

    pattern = "**/*" if recursive else "*"
    files = [
        f for f in base.glob(pattern)
        if f.is_file() and f.suffix.lower() in ALL_SUPPORTED
    ]

    if not files:
        log.warning(f"No supported files found in: {base}")
        sys.exit(0)

    log.info(f"Found {len(files)} supported file(s) in '{base}'")
    return sorted(files)


# ── CLI ───────────────────────────────────────────────────────────────────────

def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        prog="add_data.py",
        description="Ingest ANY type of file into the RAG vector database.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    p.add_argument(
        "files", nargs="*",
        help="Paths to specific files. Leave blank to open a GUI file picker.",
    )
    p.add_argument(
        "--folder", "-f", metavar="DIR",
        help="Scan a folder and ingest all supported files.",
    )
    p.add_argument(
        "--recursive", "-r", action="store_true",
        help="Scan folder recursively (use with --folder).",
    )
    p.add_argument(
        "--chunk-size", "-c", type=int, default=500, metavar="N",
        help="Maximum characters per chunk (default: 500).",
    )
    p.add_argument(
        "--overlap", "-o", type=int, default=50, metavar="N",
        help="Overlap between consecutive chunks (default: 50).",
    )
    p.add_argument(
        "--list-types", action="store_true",
        help="Print all supported file extensions and exit.",
    )
    return p


def main():
    parser = build_parser()
    args   = parser.parse_args()

    # ── --list-types ──────────────────────────────────────────────────────────
    if args.list_types:
        print("\nSupported file extensions:\n")
        for ext in sorted(ALL_SUPPORTED):
            print(f"  {ext}")
        print()
        sys.exit(0)

    # ── Resolve file list ─────────────────────────────────────────────────────
    if args.files:
        file_paths = [Path(f) for f in args.files]
    elif args.folder:
        file_paths = collect_from_folder(args.folder, recursive=args.recursive)
    else:
        # Default: open GUI file picker
        file_paths = open_file_picker()

    run_ingestion(file_paths, chunk_size=args.chunk_size, overlap=args.overlap)


if __name__ == "__main__":
    main()
